import os
import openai
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import PyPDF2


def generate_image_using_prompt(dalle_prompt):
    global images, image_url
    DALL_E_API_KEY = "YOUR_API_KEY"
    openai.api_key = DALL_E_API_KEY

    if dalle_prompt:
        try:
            response = openai.Image.create(
                prompt=dalle_prompt,
                size="512x512",  # Specify the desired image size
                num_images=1  # Number of images to generate
            )
            images = response['data']
        except Exception as e:
            print(f"Failed to generate images from DALL-E: {e}")
    else:
        print("No prompt found in the chat text.")

    # Process the generated images
    if images:
        for i, image in enumerate(images):
            image_url = image['url']
            print(f"Generated Image {i + 1}: {image_url}")
            # Process the image as needed, e.g., download, display, etc.
    else:
        print("No images generated from DALL-E.")

    # url = "https://example.com/image.jpg"
    response = requests.get(image_url)
    with open("template.jpg", "wb") as f:
        f.write(response.content)


def generate_background_image(title):
    openai.api_key = "YOUR_API_KEY"
    system_input = """
  You're a bot that helps user to generate a suitable prompt for DALL-E to generate image that can be used as a background for powerpoint presentation. It is your duty to make sure, the prompt is generated in such a way that the generated image from DALL-E doesn't contain any type of text or logo, it should be a high quality image. 
  """
    user_input = "Generate a prompt that generates a bright and light coloured image that can be used as a background for my power point presentation. Make sure the images doesn't contain any text or logo. Don't used dark colours. "

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",  # The deployment name you chose when you deployed the ChatGPT or GPT-4 model.
        messages=[
            {"role": "system", "content": system_input},
            {"role": "user", "content": user_input},
        ]

    )
    print("Title: ", title)
    print("Prompt for DALL-E: ", response['choices'][0]['message']['content'])
    prompt = response['choices'][0]['message']['content']
    generate_image_using_prompt(prompt)


def generate_content(topic):
    openai.api_key = "YOUR_API_KEY"
    system_input = """
  You're a bot that helps in generating the content for the presentation out of a given essay in the following structure
  Presentation title
  Slide-1 Title: Slide 1 Title
    - Slide 1 Content 1
    - Slide 1 Content 2
  Slide-2 Title: Slide 2 Title
    - Slide 2 Content 1
    - Slide 2 Content 2 
  ...
  Last Slide: Last Slide Title
    - Last Slide Content1
    - Last Slide Content2

  Please include "Slide-1 Title: " before adding slide title to know what's the title. Also make sure to add "Last Slide Title: " for the last slide. Also make sure you use hyphen(-) to show the slide 1 content1. Don't add sure and okay texts, just start with the main format. Also if possible, don't include "Presentation Title: ", just give the title it's understandable.
  Strictly follow the above structure no matter what.
  """
    user_input = "Give me content for a presentation based on the given essay : /n {}".format(topic)

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",  # The deployment name you chose when you deployed the ChatGPT or GPT-4 model.
        messages=[
            {"role": "system", "content": system_input},
            {"role": "user", "content": user_input},
        ]
    )

    pptx_content = response['choices'][0]['message']['content'].split("\n")
    # print("Response from the ChatGPT: \n \n", pptx_content)
    with open('pptx_content.txt', 'w') as file:
        for line in pptx_content:
            file.write(line)
            file.write("\n")


def add_background_image(slide, image_path, presentation):
    left = top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top,
                                   width=presentation.slide_width,
                                   height=presentation.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_content_to_presentation(txt_file, background_image_path, out_file_name):
    # Read the text from the input file
    with open(txt_file, "r") as file:
        lines = file.readlines()

    # Create a new presentation
    presentation = Presentation()

    # Define the slide layouts
    title_slide_layout = presentation.slide_layouts[0]  # 'Title Slide' layout
    title_and_content_layout = presentation.slide_layouts[1]  # 'Title and Content' layout

    # Add the presentation title slide
    presentation_title_slide = presentation.slides.add_slide(title_slide_layout)
    presentation_title = presentation_title_slide.shapes.title  ## Accessing the title box.
    presentation_title.text = lines.pop(0).strip()  ## Adding text in the title box.
    add_background_image(presentation_title_slide, background_image_path, presentation)

    # Loop through the lines in the input file
    slide_title = None
    slide_content = []

    for line in lines:
        if line != '\n':
            line = line.strip()
            if not line.startswith("-"):
                if slide_title is not None:
                    # Add the previous slide to the presentation
                    slide = presentation.slides.add_slide(title_and_content_layout)
                    slide_title_shape = slide.placeholders[0]
                    slide_title_shape.text = slide_title
                    slide_title_shape.text_frame.paragraphs[0].font.size = Pt(30)

                    content_placeholder = slide.placeholders[1]
                    content_text_frame = content_placeholder.text_frame  # A `TextFrame` object represents the text container associated with a shape
                    for content_line in slide_content:
                        p = content_text_frame.add_paragraph()
                        p.text = content_line
                        p.font.size = Pt(22)
                        p.space_before = Inches(0.1)

                # Start a new slide
                slide_title = line.strip().split(": ")[1]
                slide_content = []
            else:
                slide_content.append(line[1:].strip())  ## removing "-" (hyphen) from each line.

    # Add the last slide to the presentation
    slide = presentation.slides.add_slide(title_and_content_layout)
    slide_title_shape = slide.shapes.title
    slide_title_shape.text = slide_title.strip()
    slide_title_shape.text_frame.paragraphs[0].font.size = Pt(30)

    content_placeholder = slide.placeholders[1]
    content_text_frame = content_placeholder.text_frame
    for content_line in slide_content:
        p = content_text_frame.add_paragraph()
        p.text = content_line
        p.font.size = Pt(22)
        p.space_before = Inches(0.1)

    ## Add the background image to each content slide
    for index in range(1, len(presentation.slides)):
        slide = presentation.slides[index]
        add_background_image(slide, background_image_path, presentation)

    # Save the presentation to a file
    presentation.save(out_file_name)


def generate_presentation_using_topic(topic, output_file_name="presentation_using_ai.pptx"):
    openai.api_key = "YOUR_API_KEY"
    system_input = """
  You're a bot that helps in generating the content for the presentation on the given topic in the following structure
  Presentation title
  Slide-1 Title: Slide Title
    - Slide 1 Content 1
    - Slide 1 Content 2
  Slide-2 Title: Slide Title
    - Slide 2 Content 1
    - Slide 2 Content 2 
  ...
  Last Slide: Last Slide Title
    - Last Slide Content1
    - Last Slide Content2

  Please include "Slide-1 Title: " before adding slide title to know what's the title. Also make sure to add "Last Slide Title: " for the last slide. Also make sure you use hyphen(-) to show the slide 1 content1. Don't add sure, thank you and okay texts, just start with the main format. Also if possible, don't include "Presentation Title: ", just give the title it's understandable.
  Strictly follow the above structure no matter what.
  """
    user_input = "Give me content for a presentation based on the given essay : /n {}".format(topic)

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",  # The deployment name you chose when you deployed the ChatGPT or GPT-4 model.
        messages=[
            {"role": "system", "content": system_input},
            {"role": "user", "content": user_input},
        ]
    )

    pptx_content = response['choices'][0]['message']['content'].split("\n")
    topic = pptx_content[0]
    # print("Response from the ChatGPT: \n \n", pptx_content)
    with open('pptx_content.txt', 'w') as file:
        for line in pptx_content:
            file.write(line)
            file.write("\n")

    generate_background_image(topic)
    background_image_path = "template.jpg"
    add_content_to_presentation("pptx_content.txt", background_image_path, out_file_name=output_file_name)


def pdf_ext(file):
    global lst
    pdfFileObj = open(file, "rb")
    pdfReader = PyPDF2.PdfReader(pdfFileObj)
    l = []
    for i in range(len(pdfReader.pages)):
        pageObj = pdfReader.pages[i]
        if (pageObj.extract_text()) != "":
            l.append(pageObj.extract_text())
            lst = " ".join(l)
    pdfFileObj.close()
    return lst


input_text = pdf_ext("universal_declaration_of_human_rights.pdf")  ## Add pdf path
generate_presentation_using_topic(input_text)